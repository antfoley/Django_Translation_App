import subprocess
import csv
from io import BytesIO, TextIOWrapper
import docx
from django.conf import settings
from django.db import models
from django.core.files.base import ContentFile
import lxml.etree as ET
import chardet
import openpyxl
import polib
from pptx import Presentation
from google.cloud import translate_v2 as translate
from google.cloud.translate_v3 import TranslationServiceClient
from google.oauth2 import service_account
from PyPDF2 import PdfReader, PdfWriter

# try:
#     credentials = service_account.Credentials.from_service_account_file(
#         settings.GOOGLE_CLOUD_CREDENTIALS_PATH)
# except Exception as e:
credentials = service_account.Credentials.from_service_account_file(
    'translation_app/booming-post-404017-49309d69296e.json')
translateV3_client = TranslationServiceClient(credentials=credentials)
translateV2_client = translate.Client(credentials=credentials)
PARENT = f"projects/{settings.GOOGLE_PROJECT_ID}/locations/global"

def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        detector = chardet.UniversalDetector()
        for line in file:
            detector.feed(line)
            if detector.done:
                break
        detector.close()
    return detector.result['encoding']

def translate_string(target, text, source_lang):
    result = translateV2_client.translate(text,
                                          source_language=source_lang,
                                          target_language=target,
                                          format_='text')
    return str(result["translatedText"])

def translate_csv(self):
    encoding = detect_encoding(self.originalFile.path)

    original_texts = []
    rows = []
    with open(self.originalFile.path, "r", encoding=encoding) as input_file:
        reader_obj = csv.reader(input_file)
        for row in reader_obj:
            original_texts.append(row[1])
            rows.append(row)

    if self.apiVersion == 'v2':
        translated_texts = [translate_string(self.desiredLang,
                                             text,
                                             self.originalLang)
                                             for text in original_texts]
    else:
        response = translateV3_client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translated_texts = [translation.translated_text for translation in response.translations]

    output_file_content = BytesIO()
    text_wrapper = TextIOWrapper(output_file_content, encoding='utf-8', newline='')
    writer = csv.writer(text_wrapper)
    for i, row in enumerate(rows):
        row[2] = translated_texts[i]
        writer.writerow(row)

    text_wrapper.flush()
    output_file_content.seek(0)
    self.translated_file.save(f'translated_csv_file_{self.title}.csv',
                             ContentFile(output_file_content.getvalue()))
    text_wrapper.close()

def translate_xlsx(self):
    workbook = openpyxl.load_workbook(self.originalFile)
    sheet = workbook.active
    original_texts = []
    for row in sheet.iter_rows():
        for cell in row:
            if cell.value:
                original_texts.append(str(cell.value))

    if self.apiVersion == 'v2':
        translated_texts = [translate_string(self.desiredLang,
                                             text,
                                             self.originalLang)
                                             for text in original_texts]
    else:
        response = translateV3_client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translated_texts = [translation.translated_text for translation in response.translations]

    translated_text_index = 0
    for row in sheet.iter_rows():
        for cell in row:
            if cell.value:
                cell.value = translated_texts[translated_text_index]
                translated_text_index += 1

    output_file_content = BytesIO()
    workbook.save(output_file_content)
    output_file_content.seek(0)
    self.translated_file.save(f'translated_xlsx_file_{self.title}.xlsx',
                             ContentFile(output_file_content.getvalue()))

def translate_po(self):
    po_file = polib.pofile(self.originalFile.path)
    original_texts = [entry.msgid for entry in po_file if entry.msgid]

    if self.apiVersion == 'v2':
        translated_texts = [translate_string(self.desiredLang,
                                             text,
                                             self.originalLang)
                                             for text in original_texts]
    else:
        response = translateV3_client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translated_texts = [translation.translated_text for translation in response.translations]

    for i, entry in enumerate(po_file):
        if entry.msgid:
            entry.msgstr = translated_texts[i]

    output_file_content = BytesIO()
    po_file.save(output_file_content)
    output_file_content.seek(0)
    self.translated_file.save(f'translated_po_file_{self.title}.po',
                             ContentFile(output_file_content.getvalue()))

def translate_docx(self):
    doc = docx.Document(self.originalFile.path)

    original_texts = []
    translated_texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            original_texts.append(para.text)

    if self.apiVersion == 'v2':
        translated_texts = [translate_string(self.desiredLang,
                                             text,
                                             self.originalLang)
                                             for text in original_texts]
    else:
        client = translateV3_client

        response = client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translated_texts = [translation.translated_text for translation in response.translations]

    for i, para in enumerate(doc.paragraphs):
        if i < len(translated_texts):
            for run in para.runs:
                run.clear()

            translated_run = para.add_run(translated_texts[i])
            if para.runs:
                translated_run.bold = para.runs[0].bold
                translated_run.italic = para.runs[0].italic
                translated_run.underline = para.runs[0].underline
                translated_run.font.size = para.runs[0].font.size
                translated_run.font.name = para.runs[0].font.name

    output_file_content = BytesIO()
    doc.save(output_file_content)
    output_file_content.seek(0)
    self.translated_file.save(f'translated_doc_file_{self.title}.docx',
                             ContentFile(output_file_content.getvalue()))

def translate_resx(self):
    tree = ET.parse(self.originalFile.path)
    root = tree.getroot()

    original_texts = []
    translations = []

    for data in root.findall('data'):
        original_text = data.find('value').text
        if original_text:
            original_texts.append(original_text)

    if self.apiVersion == 'v2':
        translations = [translate_string(self.desiredLang,
                                         text,
                                         self.originalLang)
                                         for text in original_texts]
    else:
        client = translateV3_client
        response = client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translations = [translation.translated_text for translation in response.translations]

    for i, data in enumerate(root.findall('data')):
        if i < len(translations):  # Ensure we don't exceed the list
            value_element = data.find('value')
            if value_element is not None:
                value_element.text = translations[i]

    output_file_path = f'translated_resx_file_{self.title}.resx'
    tree.write(output_file_path, encoding='utf-8', xml_declaration=True)
    with open(output_file_path, 'rb') as file:
        contents = file.read()
    self.translated_file.save(output_file_path, contents)

def translate_doc(self):
    docx_path = self.originalFile.path.replace('.doc', '.docx')
    subprocess.run(['libreoffice',
                    '--headless',
                    '--convert-to',
                    'docx',
                    self.originalFile.path,
                    '--outdir',
                    'documents/'],
                    check=False)

    self.originalFile.name = docx_path
    translate_docx(self)

def translate_ppt(self):

    presentation = Presentation(self.originalFile.path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    if original_text.strip():
                        response = translateV3_client.translate_text(
                            parent=PARENT,
                            contents=[original_text],
                            mime_type="text/plain",
                            source_language_code=self.originalLang,
                            target_language_code=self.desiredLang,
                        )
                        translated_text = response.translations[0].translated_text
                        run.text = translated_text
                        run.font.bold = run.font.bold if run.font.bold else None
                        run.font.italic = run.font.italic if run.font.italic else None
                        run.font.underline = run.font.underline if run.font.underline else None
                        run.font.size = run.font.size if run.font.size else None
                        run.font.name = run.font.name if run.font.name else None

    output = BytesIO()
    presentation.save(output)
    content = ContentFile(output.getvalue())

    translated_filename = f'translated_ppt_file_{self.title}.pptx'
    self.translated_file.save(translated_filename, content)

def translate_pptx(self):
    prs = Presentation(self.originalFile.path)

    original_texts = []
    translations = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                if text:
                    original_texts.append(text)

    if self.apiVersion == 'v2':
        translations = [translate_string(self.desiredLang,
                                         text,
                                         self.originalLang)
                                         for text in original_texts]
    else:
        response = translateV3_client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translations = [translation.translated_text for translation in response.translations]

    translation_idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if translation_idx < len(translations):
                    shape.text = translations[translation_idx]
                    translation_idx += 1

    output_file_path = f'translated_pptx_file_{self.title}.pptx'
    prs.save(output_file_path)
    with open(output_file_path, 'rb') as file:
        contents = file.read()
    self.translated_file.save(output_file_path, contents)


def translate_pdf(self):
    reader = PdfReader(self.originalFile.path)
    writer = PdfWriter()

    original_texts = []
    translations = []

    # Extract text from each page
    for page in reader.pages:
        text = page.extract_text()
        if text:
            original_texts.append(text)

    # Translate text using the selected API version
    if self.apiVersion == 'v2':
        translations = [translate_string(self.desiredLang,
                                         text,
                                         self.originalLang)
                                         for text in original_texts]
    else:
        response = translateV3_client.translate_text(
            parent=PARENT,
            contents=original_texts,
            mime_type="text/plain",
            source_language_code=self.originalLang,
            target_language_code=self.desiredLang
        )
        translations = [translation.translated_text for translation in response.translations]

    # Create a new PDF with the translated text
    for page_num in enumerate(reader.pages):
        writer.add_page(reader.pages[page_num])
        if page_num < len(translations):
            writer.pages[page_num].insert_text(translations[page_num], position=(0, 0))
    output_file_path = f'translated_pdf_file_{self.title}.pdf'
    with open(output_file_path, 'wb') as file:
        writer.write(file)
    with open(output_file_path, 'rb') as file:
        contents = file.read()

    self.translated_file.save(output_file_path, contents)

class File(models.Model):
    title = models.CharField(max_length=200)
    originalLang = models.CharField(default='en', max_length=2)
    desiredLang = models.CharField(max_length=2)
    apiVersion = models.CharField(max_length=2)
    originalFile = models.FileField(upload_to='documents/')
    translated_file = models.FileField(upload_to='translated/', blank=True, null=True)
    def save(self, *args, **kwargs):
        super().save(*args, **kwargs)
        if self.originalFile and not self.translated_file:
            self.translate()

    def translate(self):
        arr = self.originalFile.name.split('.')
        file_type = arr[-1].lower()
        if file_type == 'csv':
            self.translated_file = f'translated_csv_file_{self.title}.csv'
            translate_csv(self)
        elif file_type == 'xlsx':
            self.translated_file = f'translated_xlsx_file_{self.title}.xlsx'
            translate_xlsx(self)
        elif file_type == 'po':
            self.translated_file = f'translated_po_file_{self.title}.po'
            translate_po(self)
        elif file_type == 'docx':
            self.translated_file = f'translated_docx_file_{self.title}.docx'
            translate_docx(self)
        elif file_type == 'doc':
            self.translated_file = f'translated_doc_file_{self.title}.doc'
            translate_doc(self)
        elif file_type == 'resx':
            self.translated_file = f'translated_resx_file_{self.title}.resx'
            translate_resx(self)
        elif file_type == 'pptx':
            self.translated_file = f'translated_pptx_file_{self.title}.pptx'
            translate_pptx(self)
        elif file_type == 'ppt':
            self.translated_file = f'translated_pptx_file_{self.title}.ppt'
            translate_pptx(self)
        elif file_type == 'pdf':
            self.translated_file = f'translated_pdf_file_{self.title}.pdf'
            translate_pdf(self)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        self.save(update_fields=['translated_file'])
